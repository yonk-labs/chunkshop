from __future__ import annotations

from pathlib import Path

from chunkshop.sources.parsers.base import ParserError


class PPTXParser:
    supported_extensions = ["pptx"]

    def parse(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except (ImportError, TypeError) as exc:
            raise RuntimeError(
                "PPTX parsing requires python-pptx. Install with `pip install chunkshop[pptx]`."
            ) from exc
        try:
            prs = Presentation(str(path))
            out: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        out.append(shape.text_frame.text)
            return "\n".join(out)
        except Exception as exc:
            raise ParserError(f"failed to parse PPTX {path}: {exc}") from exc
